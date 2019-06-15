import re
import sys
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
import codecs


class WhatsAppChatParser:
    def __init__(self, chateExportFile ):
        self.quoteList = []
        self.ignoredList = []
        self.quoteIndex = 0
        self.deletedPattern()
        self.__extractQuoteList(chateExportFile)

    def deletedPattern(self):
        messageYouDeletedMsgPattern = "^\s*\[.*\] .*: You deleted this message."
        messageOtherDeletedMsgPattern = "^\s*\[.*\] .*: This message was deleted."
        messageWebsiteLinkPattern = "^\s*\[.*\] .*: [.*https://.*|.*www..*|.*.com.*]"
        self.ignoredList.append(messageYouDeletedMsgPattern)
        self.ignoredList.append(messageOtherDeletedMsgPattern)
        self.ignoredList.append(messageWebsiteLinkPattern)
        #return(self.ignoredList)


    def shouldThisBeIgnored(self, line ):

        gotMatch = False
        for regex in self.ignoredList:
            s = re.search(regex,line)
            if( s ):
                gotMatch = True
                break
        if gotMatch:
            return (True)
        return (False)


    def __extractQuoteList(self, chateExportFile ):
        fileHandler = open (chateExportFile, "r", encoding="utf8")
        timeStamp = re.compile("^\s*\[.*\]")
        with open('config.json') as config_file:
            data = json.load(config_file)
        mFrom = data['messagesFrom']
        if ( mFrom == 'All' ):
            messageStartPattern = re.compile(".*\s*\[.*\] .*: ")
        else:
            messageStartPattern = re.compile(".*\s*\[.*\] "+mFrom+": ")
        message = ""
        insideMessage = False


        while True:
            # Get next line from file
            line = fileHandler.readline()

            #senderName = re.search('"^\s*\[.*\] (.*):', line)
            #snd = senderName.group(1)
            # If line is empty then end of file reached
            if not line :
                break;

            if ( self.shouldThisBeIgnored(line) ):
                continue

            m = messageStartPattern.match(line)
            t = timeStamp.match(line)

            if ( t ) :
                if ( insideMessage) :
                    self.quoteList.append(message)
                    insideMessage = False

                if ( m ):
                    message = line[len(m.group()):]
                    insideMessage = True
            else :
                if ( insideMessage ):
                    message = message + line



        # Close Close
        fileHandler.close()
        if ( insideMessage ):
            self.quoteList.append(message)

        #print ("MessageCount = " + str(len(self.quoteList)))
        self.slideMaking(self.quoteList)

    def slideMaking(self,quoteList):
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[6]
        i=0
        string=" "
        with open('config.json') as config_file:
            data = json.load(config_file)
        #shapeType = data['shapeType']
        #fontSize = data['fontSize']
        left = data['left']
        top = data['top']
        width = data['width']
        height = data['height']
        leftFloat = float(left)
        topFloat = float(top)
        widthFloat = float(width)
        heightFloat = float(height)
        left = Inches(leftFloat)  # 0.93" centers this overall set of shapes
        top = Inches(topFloat)
        width = Inches(widthFloat)
        height = Inches(heightFloat)
        
        for i in quoteList:
                  slide = prs.slides.add_slide(title_slide_layout)
                  shapes = slide.shapes
                  shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
                  fill = shape.fill
                  fill.solid()
                  #fill.fore_color.rgb = RGBColor(0, 0, 0)
                  fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_5
                  fill.fore_color.brightness = 0.25
                  fill.transparency = 0.25
                                   
                  text_frame = shape.text_frame
                  p = text_frame.paragraphs[0]
                  run = p.add_run()
                  run.text = i
                  font = run.font
                  font.name = 'Calibri'
                  font.size = Pt(28)
                  font.bold = False
                  font.italic = None  # cause value to be inherited from theme
                  font.color.rgb = RGBColor(0, 0, 0)
                  #font.color.theme_color = MSO_THEME_COLOR.ACCENT_8
                  
        prs.save('output.pptx')



    def getNextQuote(self):
        if ( self.quoteIndex >= len(self.quoteList)):
            raise
        message = self.quoteList[self.quoteIndex]
        self.quoteIndex += 1
        return ( message )


def testThatWeCanExtractMultiLineMessage():
    test_name = "testThatWeCanExtractMultiLineMessage"
    input_file = "testcase1.txt"
    #print(input_text)
    expected_output_text = "Ignorance is bliss\n\n"
    expected_output_text += "Its so painful to be aware of negative effects of my own actions\n\n"
    expected_output_text += "Still falling victim of this uncontrolled mind\n"
    #print(expected_output_text)

    chatparser = WhatsAppChatParser(input_file)
    actual_output = chatparser.getNextQuote()
    #print(actual_output)

    if ( actual_output == expected_output_text ):
        print(test_name + ": SUCCESS")
    else:
        print(test_name + ": FAIL" + "\n\n\tEXPECTED: " + expected_output_text + "\n\n\tACTUAL: " + actual_output)



def testThatDeletedMessagesAreIgnored():
    input_file = "deletedLinesTestData.txt"
    expected_output_text = "Quote 2\n"
    test_name = "testThatDeletedMessagesAreIgnored"

    chatparser = WhatsAppChatParser(input_file)
    first_quote = chatparser.getNextQuote()
    second_quote = chatparser.getNextQuote()

    actual_output = second_quote
    if ( actual_output == expected_output_text ):
        print(test_name + ": SUCCESS")
    else:
        print(test_name + ": FAIL" + "\n\n\tEXPECTED: " + expected_output_text + "\n\n\tACTUAL: " + actual_output)


def runTestSuite():
    testThatWeCanExtractMultiLineMessage()
    testThatDeletedMessagesAreIgnored()


def main():
    input_file = sys.argv[1]
    chatparser = WhatsAppChatParser(input_file)
    quote = chatparser.getNextQuote()
    print(quote + "\n#######\n")
    quote = chatparser.getNextQuote()
    print(quote + "\n#######\n")
    quote = chatparser.getNextQuote()
    print(quote + "\n#######\n")
    quote = chatparser.getNextQuote()
    print(quote + "\n#######\n")



#runTestSuite()
main()
