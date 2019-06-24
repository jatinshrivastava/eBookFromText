import sys
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
import codecs


class SlideGenerator:
    def __init__(self, pptFileName ):
        self.pptFileName = pptFileName
        with open('config.json') as config_file:
            data = json.load(config_file)
        self.shapeType = data['shapeType']
        self.fontSize = int(data['fontSize'])
        self.left = Inches(float(data['left']))
        self.top = Inches(float(data['top']))
        self.width = Inches(float(data['width']))
        self.height = Inches(float(data['height']))
        self.shapemap = {}
        self.shapemap['rectangle'] = MSO_SHAPE.RECTANGLE
        self.shapemap['round rectangle'] = MSO_SHAPE.ROUNDED_RECTANGLE
        self.shapemap['curved ribbon'] = MSO_SHAPE.CURVED_UP_RIBBON

        self.prs = Presentation()
        self.title_slide_layout = self.prs.slide_layouts[6]

    def addSlide(self, quote ):
        slide = self.prs.slides.add_slide(self.title_slide_layout)
        shapes = slide.shapes
        shape = shapes.add_shape(self.shapemap[self.shapeType], self.left, self.top, self.width, self.height)
        fill = shape.fill
        fill.solid()
        #fill.fore_color.rgb = RGBColor(0, 0, 0)
        fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_5
        fill.fore_color.brightness = 0.25
        fill.transparency = 0.25
                       
        text_frame = shape.text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = quote
        font = run.font
        font.name = 'Calibri'
        font.size = Pt(self.fontSize)
        font.bold = False
        font.italic = None  # cause value to be inherited from theme
        font.color.rgb = RGBColor(0, 0, 0)
        #font.color.theme_color = MSO_THEME_COLOR.ACCENT_8
        
    def save(self ):
        self.prs.save(self.pptFileName)
  
