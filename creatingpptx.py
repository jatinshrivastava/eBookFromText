from pptx import Presentation #TODO

prs = Presentation() #TODO
title_slide_layout = prs.slide_layouts[0] #you can change the layout type here.
slide = prs.slides.add_slide(title_slide_layout) #TODO
title = slide.shapes.title #TODO
subtitle = slide.placeholders[1] #TODO

#testing of Accessing of file.
f1 = open("info.txt","r") # r here stands for reading of file.
print ("Name of file : ", f1.name) #acessing the file name here.  
f1.mode == 'r'
content1 = f1.read() #storing the content of the file into a variable (which is content1 here).
print(content1) #printing the contents here.
#content1 below is the variable that stores the content of text file.
title.text = content1  #this is the text inside the that will appear on slide.
subtitle.text = "python-pptx was here!" #this is the text inside the that will appear on slide.

prs.save('test_ppt.pptx') #this will affect the file name.